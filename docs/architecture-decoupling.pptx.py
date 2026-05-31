"""Generate docs/architecture-decoupling.pptx — code architecture +
rover-drone decoupling deck.

Usage:
    python3 docs/architecture-decoupling.pptx.py
    # writes docs/architecture-decoupling.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).parent / "architecture-decoupling.pptx"

# Palette
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x21, 0x6E, 0x9B)
DRONE = RGBColor(0x4A, 0x90, 0xE2)
ROVER = RGBColor(0xE2, 0x8B, 0x4A)
CONTROLLER = RGBColor(0x6F, 0x42, 0xC1)
PIPELINE = RGBColor(0x2D, 0xA0, 0x6C)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = INK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.runs[0].font.size = Pt(15)
        p2.runs[0].font.color.rgb = MUTED
        p2.runs[0].font.italic = True


def add_text_box(slide, left, top, width, height, lines, *,
                 size=14, color=INK, bold=False, italic=False, mono=False,
                 align=PP_ALIGN.LEFT, bullet=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        prefix = "• " if bullet else ""
        # python-pptx leaves p.runs empty if text is empty; use a space sentinel
        p.text = (prefix + line) if line else " "
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        if mono:
            r.font.name = "Consolas"
    return box


def add_box(slide, left, top, width, height, fill, *, text=None,
            text_color=None, text_size=14, text_bold=True,
            border=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.75)
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.text = line
            r = p.runs[0]
            r.font.size = Pt(text_size)
            r.font.bold = text_bold and i == 0
            r.font.color.rgb = text_color or BG_WHITE
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, label=None,
              label_offset_x=Inches(0), label_offset_y=Inches(-0.18)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.end_arrow_type = 2  # arrow tail
    if label:
        lab = add_text_box(
            slide,
            min(x1, x2) + (abs(x2 - x1) // 2) - Inches(0.6) + label_offset_x,
            (y1 + y2) // 2 + label_offset_y,
            Inches(1.6), Inches(0.3),
            [label], size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
        )
    return line


def set_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="drone-follow / hailo-drone-follow • v1.1 (Phase 6)"):
    add_text_box(
        slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
        [text], size=9, color=MUTED, italic=True, align=PP_ALIGN.LEFT,
    )


# ---------- Build deck ----------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# Slide 1 — Title
s = prs.slides.add_slide(blank)
set_bg(s, BG_WHITE)
add_text_box(
    s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.0),
    ["Code Architecture & Decoupling"],
    size=42, bold=True, color=INK, align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.7),
    ["drone-follow / hailo-drone-follow — v1.1 rover support"],
    size=22, color=ACCENT, align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.0),
    [
        "How the system stays robot-agnostic: one controller, two adapters,",
        "a small contract surface, zero leakage of robot knowledge into the controller.",
    ],
    size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
    ["2026-05-20"],
    size=12, color=MUTED, align=PP_ALIGN.CENTER,
)


# Slide 2 — Four-layer architecture
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Four-Layer Architecture",
          "Each layer depends only on the layer below it.")

LX, LW = Inches(2.0), Inches(9.3)
LH = Inches(0.85)
y = Inches(1.4)
layers = [
    ("Servers (UI + Telemetry)", "follow_server.py · web_server.py · openhd_bridge.py", ACCENT),
    ("Pipeline Adapter (vision)", "Hailo + GStreamer · ByteTracker · ReID manager", PIPELINE),
    ("Follow API (domain)", "controller · types (Capabilities, RobotCommand, SafetyContext) · config", CONTROLLER),
    ("Robot API (actuator boundary)", "Robot Protocol · MavsdkDroneAdapter · Ros2RoverAdapter", DRONE),
]
for title, sub, color in layers:
    add_box(s, LX, y, LW, LH, color, text=title, text_size=17)
    add_text_box(
        s, LX + LW + Inches(0.2), y + Inches(0.1), Inches(2.0), LH - Inches(0.1),
        [sub], size=11, color=MUTED, italic=True,
    )
    y += LH + Inches(0.15)

add_text_box(
    s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
    [
        "Key invariant: the lower the layer, the less it knows.",
        "follow_api/ knows nothing about drones, ROS, or MAVSDK — it just emits "
        "RobotCommand + SafetyContext.",
    ],
    size=14, color=INK, bullet=True,
)
add_footer(s)


# Slide 3 — Who knows whom
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Who Knows Whom — Dependency Direction",
          "Arrows point FROM consumer TO contract. No back-arrows.")

# Top row: pipeline + servers (consumers of controller)
add_box(s, Inches(0.7), Inches(1.5), Inches(3.5), Inches(0.9), PIPELINE,
        text="pipeline_adapter\n(Hailo + ByteTracker)", text_size=14)
add_box(s, Inches(9.1), Inches(1.5), Inches(3.5), Inches(0.9), ACCENT,
        text="servers/\n(web_server, openhd_bridge)", text_size=14)

# Middle row: controller + orchestrator
add_box(s, Inches(2.8), Inches(3.0), Inches(3.5), Inches(0.9), CONTROLLER,
        text="follow_api.controller", text_size=15)
add_box(s, Inches(7.0), Inches(3.0), Inches(3.5), Inches(0.9), CONTROLLER,
        text="robot_api.orchestrator", text_size=15)

# Boundary row: Robot Protocol
add_box(s, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.85),
        INK, text="Robot Protocol\n(robot_api/robot.py)", text_size=15)

# Bottom row: adapters
add_box(s, Inches(2.0), Inches(6.0), Inches(3.8), Inches(0.85), DRONE,
        text="MavsdkDroneAdapter", text_size=14)
add_box(s, Inches(7.5), Inches(6.0), Inches(3.8), Inches(0.85), ROVER,
        text="Ros2RoverAdapter", text_size=14)

# Arrows
add_arrow(s, Inches(2.4), Inches(2.4), Inches(4.0), Inches(3.0),
          color=PIPELINE, label="passes RobotCommand")
add_arrow(s, Inches(10.85), Inches(2.4), Inches(8.7), Inches(3.0),
          color=ACCENT, label="reads SharedState")
add_arrow(s, Inches(6.3), Inches(3.45), Inches(7.0), Inches(3.45),
          color=MUTED, label="emits cmd")
add_arrow(s, Inches(8.7), Inches(3.9), Inches(7.0), Inches(4.5),
          color=MUTED, label="depends on")
add_arrow(s, Inches(6.65), Inches(5.35), Inches(3.9), Inches(6.0),
          color=DRONE)
add_arrow(s, Inches(6.65), Inches(5.35), Inches(9.4), Inches(6.0),
          color=ROVER)

add_text_box(
    s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
    ["Critical: there is NO arrow from controller → adapter. The controller depends "
     "on Robot (interface), not MavsdkDroneAdapter or Ros2RoverAdapter."],
    size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)


# Slide 4 — Robot Protocol = decoupling boundary
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Robot Protocol — The Decoupling Boundary",
          "Five async methods + one Capabilities attr. That's the whole contract.")

add_text_box(
    s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
    ["robot_api/robot.py"],
    size=14, mono=True, color=ACCENT, bold=True,
)
code = [
    "@runtime_checkable",
    "class Robot(Protocol):",
    "    caps: Capabilities",
    "",
    "    async def connect(self) -> None: ...",
    "    async def start_session(self) -> None: ...",
    "    async def send_command(self, cmd: RobotCommand,",
    "                           safety_ctx: SafetyContext) -> None: ...",
    "    async def send_zero(self) -> None: ...",
    "    async def on_target_lost(self,",
    "                             last_detection: Optional[Detection]) -> None: ...",
    "    async def shutdown(self) -> None: ...",
]
add_text_box(
    s, Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.5),
    code, size=14, mono=True, color=INK,
)

# Right side — semantic notes
notes = [
    "send_command",
    "  Adapter overlays its robot-specific behaviors here.",
    "  MUST early-return on safety_ctx.target_lost == True.",
    "",
    "yaw_rate units",
    "  Drone: deg/s. Rover: rad/s.",
    "  NO adapter-side conversion (Q5 lock).",
    "",
    "shutdown",
    "  Idempotent; always called in orchestrator's finally.",
]
add_box(s, Inches(8.4), Inches(1.95), Inches(4.4), Inches(4.6), BG_WHITE,
        border=MUTED, text="")
y = Inches(2.1)
for line in notes:
    bold = bool(line) and not line.startswith("  ")
    italic = line.startswith("  ")
    add_text_box(
        s, Inches(8.55), y, Inches(4.1), Inches(0.35),
        [line], size=12, color=INK if bold else MUTED,
        bold=bold, italic=italic, mono=False,
    )
    y += Inches(0.36)


# Slide 5 — Data contracts
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Data Contracts Across the Boundary",
          "Three small dataclasses. No I/O, no behavior, just shape.")

# Three columns
col_w = Inches(4.05)
col_h = Inches(5.2)
xs = [Inches(0.5), Inches(4.65), Inches(8.8)]

# Capabilities
add_box(s, xs[0], Inches(1.4), col_w, Inches(0.7), CONTROLLER,
        text="Capabilities", text_size=18)
add_text_box(
    s, xs[0] + Inches(0.1), Inches(2.2), col_w - Inches(0.2), col_h,
    [
        "Axes-only, by design",
        "",
        "axes: frozenset[Axis]",
        "  FORWARD · YAW · ALTITUDE",
        "",
        "yaw_unit: \"deg/s\" | \"rad/s\"",
        "",
        "NO behavioral flags.",
        "No bottom_edge_policy.",
        "No yaw_spin_on_loss.",
        "",
        "→ Behaviors live in the",
        "   adapter, never gated by a",
        "   Capabilities flag.",
    ],
    size=13, color=INK,
)

# RobotCommand
add_box(s, xs[1], Inches(1.4), col_w, Inches(0.7), ACCENT,
        text="RobotCommand", text_size=18)
add_text_box(
    s, xs[1] + Inches(0.1), Inches(2.2), col_w - Inches(0.2), col_h,
    [
        "Per-tick actuator command",
        "",
        "forward_m_s: float",
        "  body x velocity, m/s",
        "",
        "yaw_rate: float",
        "  rotation about body z",
        "  units: caps.yaw_unit",
        "",
        "down_m_s: float",
        "  body z, m/s (down +)",
        "  only when Axis.ALTITUDE",
        "",
        "→ Controller writes only what",
        "   caps.axes permits.",
    ],
    size=13, color=INK,
)

# SafetyContext
add_box(s, xs[2], Inches(1.4), col_w, Inches(0.7), DRONE,
        text="SafetyContext", text_size=18)
add_text_box(
    s, xs[2] + Inches(0.1), Inches(2.2), col_w - Inches(0.2), col_h,
    [
        "Per-tick safety signals",
        "",
        "bbox_bottom_normalized",
        "  drone retreat-from-tilt",
        "",
        "bbox_bottom_norm (Optional)",
        "  rover slow-down (RINT-02)",
        "",
        "bbox_size_normalized",
        "target_lost: bool",
        "last_target_x: Optional[float]",
        "",
        "→ Adapter MUST early-return",
        "   if target_lost is True.",
    ],
    size=13, color=INK,
)


# Slide 6 — Drone adapter
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "MavsdkDroneAdapter — What the Drone Adapter Owns",
          "Adapter-local behaviors that never leak into the controller.")

add_box(s, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.7), DRONE,
        text="caps", text_size=14)
add_text_box(
    s, Inches(3.2), Inches(1.45), Inches(9.0), Inches(0.65),
    ["axes = {FORWARD, YAW, ALTITUDE} · yaw_unit = \"deg/s\""],
    size=14, mono=True, color=INK,
)

items = [
    ("Wire", "MAVSDK over UDP/serial · OFFBOARD setpoint stream · VelocityBodyYawspeed"),
    ("Altitude P-loop", "fixed-altitude hold via target_altitude (controller is altitude-blind)"),
    ("Retreat-from-tilt", "reads safety_ctx.bbox_bottom_normalized, eases back when bbox sinks"),
    ("Slew-rate cap", "_apply_smoothing limits forward-axis acceleration (control_loop_hz aware)"),
    ("Yaw units", "controller emits deg/s — drone publishes deg/s. NO conversion."),
    ("Takeoff / land", "owned by adapter when --takeoff-landing; otherwise pilot owns mode switch"),
    ("Mission watchdog", "auto-land on duration expiry (with takeoff-landing); restart otherwise"),
]
y = Inches(2.3)
for label, body in items:
    add_box(s, Inches(0.5), y, Inches(2.5), Inches(0.55), DRONE,
            text=label, text_size=13)
    add_text_box(
        s, Inches(3.2), y + Inches(0.08), Inches(9.6), Inches(0.5),
        [body], size=13, color=INK,
    )
    y += Inches(0.65)


# Slide 7 — Rover adapter
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Ros2RoverAdapter — What the Rover Adapter Owns",
          "Same Robot interface — but a completely different set of robot-specific behaviors.")

add_box(s, Inches(0.5), Inches(1.4), Inches(2.5), Inches(0.7), ROVER,
        text="caps", text_size=14)
add_text_box(
    s, Inches(3.2), Inches(1.45), Inches(9.0), Inches(0.65),
    ["axes = {FORWARD, YAW} · yaw_unit = \"rad/s\"   (NO ALTITUDE)"],
    size=14, mono=True, color=INK,
)

items = [
    ("Wire", "rclpy node + /cmd_vel Twist publisher (ROS 2 Humble · DDS)"),
    ("NED→ENU yaw flip", "controller emits NED CW+ → adapter publishes angular.z = -yaw (REP-103 ENU CCW+)"),
    ("Bottom-edge slow-stop", "RINT-02 · twist.linear.x = 0 when safety_ctx.bbox_bottom_norm ≥ 0.85"),
    ("Yaw preserved on stop", "twist.angular.z stays at cmd.yaw_rate so the rover can recenter"),
    ("EMA smoothing", "per-axis low-pass on yaw + forward (config.smooth_yaw / smooth_forward)"),
    ("Yaw-spin search on loss", "on_target_lost yaws toward last-seen side at config.search_yawspeed_slow"),
    ("Yaw units", "controller emits rad/s — rover publishes rad/s. NO conversion."),
    ("No altitude", "Axis.ALTITUDE not in caps → controller never emits down_m_s"),
    ("SIGINT clean stop", "RINT-06 · publisher cleared in shutdown(); destroy_node → try_shutdown"),
]
y = Inches(2.3)
for label, body in items:
    add_box(s, Inches(0.5), y, Inches(2.5), Inches(0.55), ROVER,
            text=label, text_size=13)
    add_text_box(
        s, Inches(3.2), y + Inches(0.08), Inches(9.6), Inches(0.5),
        [body], size=13, color=INK,
    )
    y += Inches(0.65)


# Slide 8 — Controller stays robot-agnostic
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "How the Controller Stays Robot-Agnostic",
          "What is — and isn't — in follow_api/controller.py.")

add_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.7),
        PIPELINE, text="✓ Controller KNOWS", text_size=16)
add_text_box(
    s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(4.5),
    [
        "Per-tick bbox geometry",
        "  cx, cy, bbox_height in normalized coords",
        "",
        "Yaw-axis P controller",
        "  emits yaw_rate in caps.yaw_unit",
        "",
        "Forward-axis distance keep",
        "  shrinks/grows command to hit target_bbox_height",
        "",
        "Yaw smoothing (EMA + dead zone)",
        "",
        "Search-mode escalation",
        "  on detection_timeout_s → on_target_lost",
        "",
        "What axes are allowed",
        "  reads caps.axes to know which channels to write",
    ],
    size=13, color=INK,
)

add_box(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.7),
        DRONE, text="✗ Controller does NOT know", text_size=16)
add_text_box(
    s, Inches(6.9), Inches(2.2), Inches(5.9), Inches(4.5),
    [
        "What kind of robot it's driving",
        "  no isinstance checks on adapter",
        "",
        "Wire protocol",
        "  no MAVSDK · no rclpy · no ROS · no DDS",
        "",
        "Altitude",
        "  drone owns the alt-hold loop in mavsdk_drone.py",
        "",
        "Bottom-edge slow-down",
        "  rover-specific; lives in ros2_rover.py",
        "",
        "Takeoff / land",
        "  drone-only concept; controller never mentions it",
        "",
        "Unit conversions",
        "  the controller emits in caps.yaw_unit, period",
    ],
    size=13, color=INK,
)


# Slide 9 — Worked example: RINT-02
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Worked Example: RINT-02 Rover Bottom-Edge Slow-Down",
          "Why the threshold lives in the adapter, not Capabilities, not the controller.")

# Left: the change site
add_text_box(
    s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.4),
    ["robot_follow/robot_api/adapters/ros2_rover.py"],
    size=13, mono=True, color=ACCENT, bold=True,
)
code = [
    "ROVER_BOTTOM_STOP_THRESHOLD: float = 0.85",
    "",
    "async def send_command(self, cmd, safety_ctx):",
    "    if safety_ctx.target_lost:",
    "        return",
    "    if self._publisher is None:",
    "        return",
    "    twist = self._Twist()",
    "    twist.linear.x  = cmd.forward_m_s",
    "    twist.angular.z = cmd.yaw_rate",
    "    if (safety_ctx.bbox_bottom_norm is not None",
    "            and safety_ctx.bbox_bottom_norm",
    "                >= ROVER_BOTTOM_STOP_THRESHOLD):",
    "        twist.linear.x = 0.0     # stop forward",
    "        # yaw preserved → rover can still",
    "        # rotate to recenter the actor",
    "    self._publisher.publish(twist)",
]
add_text_box(
    s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.2),
    code, size=12, mono=True, color=INK,
)

# Right: rationale
add_text_box(
    s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.4),
    ["Why here — not in Capabilities or controller"],
    size=14, color=INK, bold=True,
)
add_text_box(
    s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(5.0),
    [
        "Robot-specific physics",
        "  A rover at ground level hits things. A drone at 3m altitude doesn't.",
        "",
        "Named constant in ONE site",
        "  Greppable: `ROVER_BOTTOM_STOP_THRESHOLD`. Tuning is a single edit.",
        "",
        "Capabilities stays axes-only",
        "  No bottom_edge_policy flag would have smuggled rover knowledge",
        "  back into the controller.",
        "",
        "Drone adapter unchanged",
        "  mavsdk_drone.py was byte-identical across all 7 Phase-6 plans.",
        "  The drone's own retreat-from-tilt continues to read",
        "  bbox_bottom_normalized — a separate field, separate adapter.",
        "",
        "Twist-only override",
        "  cmd (the RobotCommand) is never mutated — only the about-to-",
        "  publish twist. No surprising upstream side-effects.",
    ],
    size=13, color=INK,
)


# Slide 10 — Decoupling principles summary
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Five Decoupling Principles",
          "What kept the controller swap-compatible from drone to rover.")

principles = [
    ("1", "Protocol over concrete classes",
     "Robot is a typing.Protocol. Controller depends on the interface, never on "
     "MavsdkDroneAdapter or Ros2RoverAdapter."),
    ("2", "Capabilities = axes only",
     "frozenset[Axis] + yaw_unit. No behavioral flags. Behaviors belong inside "
     "the adapter, not in a shared dataclass."),
    ("3", "Adapter owns unit semantics",
     "Drone emits deg/s, rover emits rad/s. Controller writes in caps.yaw_unit; "
     "no adapter does any conversion."),
    ("4", "Data flows through small dataclasses",
     "RobotCommand + SafetyContext are the only types crossing the boundary. "
     "Both are pure data — no I/O, no behavior."),
    ("5", "Extension is additive",
     "Adding the rover required: 1 new adapter, 1 enum value, 1 optional "
     "SafetyContext field. Zero changes to controller logic."),
]
y = Inches(1.4)
for n, title, body in principles:
    add_box(s, Inches(0.5), y, Inches(0.7), Inches(0.95), ACCENT,
            text=n, text_size=22)
    add_text_box(
        s, Inches(1.4), y, Inches(11.5), Inches(0.4),
        [title], size=16, bold=True, color=INK,
    )
    add_text_box(
        s, Inches(1.4), y + Inches(0.4), Inches(11.5), Inches(0.55),
        [body], size=12, color=MUTED,
    )
    y += Inches(1.08)


# Slide 11 — Forensics: the JSON-loading trap
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Forensics — argparse defaults silently shadowed JSON config",
          "Bug that made the rover behave with drone gains regardless of --config.")

add_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.6),
        RGBColor(0xCC, 0x44, 0x44), text="Before (broken)", text_size=14)
add_text_box(
    s, Inches(0.5), Inches(2.05), Inches(6.0), Inches(2.4),
    [
        "# config.add_args",
        "parser.add_argument(",
        "    '--yaw-gain', dest='kp_yaw',",
        "    type=float,",
        "    default=defaults.kp_yaw,   # ← 4.0 (drone)",
        ")",
        "",
        "# config.from_args:",
        "kp_yaw = _arg('kp_yaw', default=defaults.kp_yaw)",
        "# args.kp_yaw is ALWAYS 4.0 even with no flag",
        "# _arg returns it; JSON's 0.05 is discarded.",
    ],
    size=12, mono=True, color=INK,
)

add_box(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.6),
        PIPELINE, text="After (fixed in 9b2e978)", text_size=14)
add_text_box(
    s, Inches(6.8), Inches(2.05), Inches(6.0), Inches(2.4),
    [
        "# config.add_args",
        "parser.add_argument(",
        "    '--yaw-gain', dest='kp_yaw',",
        "    type=float,",
        "    default=None,              # ← key change",
        ")",
        "",
        "# config.from_args — unchanged:",
        "kp_yaw = _arg('kp_yaw', default=defaults.kp_yaw)",
        "# args.kp_yaw is None ⇒ _arg falls through",
        "# to defaults.kp_yaw = JSON's 0.05.",
    ],
    size=12, mono=True, color=INK,
)

add_text_box(
    s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
    ["Observed (rover sim, /tmp/rover.log):"],
    size=14, bold=True, color=INK,
)
add_text_box(
    s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.45),
    ["yaw raw = -20.95 at cx = 0.084   matches  kp_yaw=4.0 × √27.5°"],
    size=12, mono=True, color=INK,
)
add_text_box(
    s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.45),
    ["Expected with JSON 0.05:   yaw raw = -0.26    (80× smaller)"],
    size=12, mono=True, color=MUTED,
)

add_text_box(
    s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(1.0),
    [
        "Second bug, same fix commit: from_args constructs cls(...) with an explicit kwarg list.",
        "max_yawspeed wasn't in the list → JSON value silently fell through to the dataclass default.",
        "Rover JSON's max_yawspeed = 0.8 rad/s was discarded → 90.0 (drone deg/s) applied as rad/s.",
    ],
    size=12, color=INK,
)

add_text_box(
    s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.5),
    ["Regression test: test_from_args_does_not_override_json_with_cli_defaults"],
    size=12, italic=True, color=MUTED,
)


# Slide 12 — Recap
s = prs.slides.add_slide(blank)
set_bg(s, BG_WHITE)
add_text_box(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.0),
    ["One controller. Two adapters. One small contract."],
    size=32, bold=True, color=INK, align=PP_ALIGN.CENTER,
)
add_text_box(
    s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(2.5),
    [
        "follow_api/controller.py: never mentions \"drone\" or \"rover\".",
        "mavsdk_drone.py: never mentions ros2_rover. ros2_rover.py: never mentions mavsdk.",
        "Capabilities never grew a behavioral flag.",
        "",
        "v1.1 added rover support by adding code, not by editing the controller.",
    ],
    size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)
add_footer(s)


prs.save(str(OUT_PATH))
print(f"wrote {OUT_PATH}")
