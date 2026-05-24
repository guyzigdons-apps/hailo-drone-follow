"""Generate docs/v1.1-rover-support.pptx — a short focused deck on the
feature/rover-support branch: what changed, what was decoupled, what
the new orchestrator is for, and which fixes this session introduced.

Usage:
    python3 docs/v1.1-rover-support.pptx.py
    # writes docs/v1.1-rover-support.pptx (gitignored via *.pptx)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_PATH = Path(__file__).parent / "v1.1-rover-support.pptx"

# Palette — shared with architecture-decoupling.pptx for visual continuity.
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
SOFT = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0x21, 0x6E, 0x9B)
DRONE = RGBColor(0x4A, 0x90, 0xE2)
ROVER = RGBColor(0xE2, 0x8B, 0x4A)
CONTROLLER = RGBColor(0x6F, 0x42, 0xC1)
PIPELINE = RGBColor(0x2D, 0xA0, 0x6C)
WARN = RGBColor(0xCC, 0x44, 0x44)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ── helpers (same shapes used in architecture-decoupling.pptx) ──────────

def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = INK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.runs[0].font.size = Pt(15)
        p2.runs[0].font.color.rgb = MUTED
        p2.runs[0].font.italic = True


def add_text(slide, left, top, width, height, lines, *,
             size=14, color=INK, bold=False, italic=False, mono=False,
             align=PP_ALIGN.LEFT, bullet=False, line_spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        prefix = "• " if bullet else ""
        p.text = (prefix + line) if line else " "
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        if mono:
            r.font.name = "Consolas"
    return box


def add_box(slide, left, top, width, height, fill, *, text=None,
            text_color=None, text_size=14, text_bold=True,
            border=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(0.75)
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.text = line
            r = p.runs[0]
            r.font.size = Pt(text_size)
            r.font.bold = text_bold and i == 0
            r.font.color.rgb = text_color or BG_WHITE
    return shp


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_footer(slide):
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3),
        ["feature/rover-support · v1.1"],
        size=10, color=SOFT, italic=True, align=PP_ALIGN.RIGHT,
    )


# ── deck ────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]


# Slide 1 — Cover
s = prs.slides.add_slide(blank)
set_bg(s, BG_WHITE)
add_text(
    s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.0),
    ["Drone Follow → Robot Follow"],
    size=42, bold=True, color=INK, align=PP_ALIGN.CENTER,
)
add_text(
    s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.7),
    ["v1.1 · the rover-support branch"],
    size=22, color=ACCENT, italic=True, align=PP_ALIGN.CENTER,
)
add_text(
    s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.6),
    [
        "200 commits ahead of main · adds a pluggable robot-adapter layer",
        "so one controller can drive a PX4 drone or a ROS 2 skid-steer rover",
        "without knowing which one it is.",
    ],
    size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)


# Slide 2 — The big picture
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "What this branch does",
          "Four interlocking changes; everything else is in service of them.")

steps = [
    ("1", "Renamed", "drone_follow → robot_follow", DRONE,
     "Python package + console-script. `drone-follow` alias preserved permanently so the boot service + field deployments don't break."),
    ("2", "Extracted", "Robot protocol + Capabilities", CONTROLLER,
     "Tiny duck-typed interface that the controller and orchestrator depend on. The drone adapter and rover adapter implement it. Axes-only — no behavioral flags."),
    ("3", "Added", "Ros2RoverAdapter", ROVER,
     "geometry_msgs/Twist on /cmd_vel · rad/s yaw · NED→ENU sign flip · bottom-edge slow-stop · yaw-spin search on target loss."),
    ("4", "Added", "Generic orchestrator", PIPELINE,
     "robot_api/orchestrator.py — one per-tick control loop both adapters share. Reads detection state, calls send_command or on_target_lost."),
]
y = Inches(1.45)
for n, verb, title, color, body in steps:
    add_box(s, Inches(0.5), y, Inches(0.65), Inches(1.05), color,
            text=n, text_size=26)
    add_text(
        s, Inches(1.35), y, Inches(2.0), Inches(0.45),
        [verb], size=12, color=MUTED, italic=True,
    )
    add_text(
        s, Inches(1.35), y + Inches(0.35), Inches(11.5), Inches(0.5),
        [title], size=18, bold=True, color=INK,
    )
    add_text(
        s, Inches(1.35), y + Inches(0.78), Inches(11.5), Inches(0.45),
        [body], size=12, color=MUTED,
    )
    y += Inches(1.25)


# Slide 3 — The new layer
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "The new robot_api/ layer",
          "All adapter-shaped concerns now live behind one boundary.")

# Two-column layout: tree on left, role text on right.
tree_lines = [
    "robot_follow/",
    "├── follow_api/",
    "│     ├── controller.py     ◀ pure P-controller — sees only bbox geometry",
    "│     ├── config.py",
    "│     └── types.py          ◀ RobotCommand · SafetyContext · Capabilities · Axis",
    "│",
    "├── robot_api/              ◀ NEW",
    "│     ├── robot.py            ▸ Robot protocol",
    "│     ├── orchestrator.py     ▸ generic per-tick loop  ★",
    "│     └── adapters/",
    "│           ├── mavsdk_drone.py    ▸ PX4 offboard + telemetry",
    "│           └── ros2_rover.py      ▸ ROS 2 Twist + diff-drive",
    "│",
    "├── pipeline_adapter/         (Hailo · ByteTracker · ReID)",
    "├── servers/                  (follow API · web UI · OpenHD bridge)",
    "└── robot_follow_app.py       (composition root + --robot {drone,rover})",
]
add_text(
    s, Inches(0.5), Inches(1.4), Inches(8.3), Inches(5.0),
    tree_lines, size=12, mono=True, color=INK, line_spacing=1.15,
)

add_box(s, Inches(9.1), Inches(1.4), Inches(3.7), Inches(0.55), CONTROLLER,
        text="The keystone", text_size=14)
add_text(
    s, Inches(9.1), Inches(2.05), Inches(3.7), Inches(4.5),
    [
        "robot_api/ is a thin layer:",
        "",
        "▸ robot.py defines what every",
        "  adapter must implement",
        "  (connect, send_command,",
        "  on_target_lost, send_zero,",
        "  shutdown, plus a caps attr).",
        "",
        "▸ orchestrator.py runs the",
        "  loop. Robot-agnostic.",
        "",
        "▸ adapters/ holds the only",
        "  code that talks to MAVSDK",
        "  or rclpy. Each adapter is",
        "  swappable in isolation.",
    ],
    size=12, color=INK,
)


# Slide 4 — Decoupled vs coupled
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Decoupled / Coupled",
          "What each side knows — and what crosses the boundary.")

# Decoupled column
add_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.6),
        PIPELINE, text="✂  Decoupled  (adapter-side)", text_size=14)
decoupled = [
    "MAVSDK · rclpy · ROS 2 imports",
    "deg/s vs rad/s yaw units",
    "NED ↔ ENU frame translation",
    "Takeoff / landing semantics",
    "Altitude-hold P loop  (drone only)",
    "Tilt-induced retreat  (drone only)",
    "Bottom-edge slow-stop  (rover only)",
    "Search-mode yaw behaviour on target loss",
    "Per-axis EMA smoothing",
    "Lifecycle shutdown details  (SIGINT, executor thread, mavsdk_server)",
]
add_text(
    s, Inches(0.5), Inches(2.05), Inches(6.0), Inches(5.0),
    decoupled, size=13, color=INK, bullet=True, line_spacing=1.35,
)

# Coupled column
add_box(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.6),
        ACCENT, text="🔗  Coupled  (shared contract)", text_size=14)
coupled = [
    "Robot protocol  (5 methods + caps)",
    "RobotCommand  (forward_m_s, yaw_rate, down_m_s)",
    "SafetyContext  (bbox geometry, target_lost)",
    "Capabilities  (axes + yaw_unit only — no policy)",
    "ControllerConfig  (all gains; per-axis tunables)",
    "follow_api/controller.compute  (pure P math)",
    "robot_api/orchestrator.run_robot_loop",
    "FollowTargetState + SharedDetectionState",
    "web UI + OpenHD bridge",
    "—  the boundary is small on purpose  —",
]
add_text(
    s, Inches(6.8), Inches(2.05), Inches(6.0), Inches(5.0),
    coupled, size=13, color=INK, bullet=True, line_spacing=1.35,
)
add_footer(s)


# Slide 5 — The orchestrator
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "The Orchestrator  —  one loop for any robot",
          "robot_api/orchestrator.run_robot_loop — robot-agnostic per-tick control.")

# Lifecycle row at the top
add_box(s, Inches(0.5), Inches(1.4), Inches(2.8), Inches(0.65),
        ACCENT, text="connect()", text_size=14)
add_text(s, Inches(3.4), Inches(1.5), Inches(0.3), Inches(0.45), ["▸"], size=22, color=MUTED)
add_box(s, Inches(3.8), Inches(1.4), Inches(2.8), Inches(0.65),
        ACCENT, text="start_session()", text_size=14)
add_text(s, Inches(6.7), Inches(1.5), Inches(0.3), Inches(0.45), ["▸"], size=22, color=MUTED)
add_box(s, Inches(7.1), Inches(1.4), Inches(2.8), Inches(0.65),
        ACCENT, text="per-tick loop", text_size=14)
add_text(s, Inches(10.0), Inches(1.5), Inches(0.3), Inches(0.45), ["▸"], size=22, color=MUTED)
add_box(s, Inches(10.4), Inches(1.4), Inches(2.4), Inches(0.65),
        ACCENT, text="send_zero · shutdown", text_size=14)

# State machine — three branches under "per-tick loop"
add_text(
    s, Inches(0.5), Inches(2.35), Inches(12.3), Inches(0.4),
    ["Per-tick state machine"], size=14, bold=True, color=INK,
)

branches = [
    ("Detection present",
     "compute(det, caps, config) → send_command(cmd, safety_ctx)",
     "Fresh setpoint every tick. Adapter applies its own smoothing / safety.",
     PIPELINE),
    ("Absent · within search_enter_delay_s",
     "re-send last_cmd with safety_ctx from last_detection",
     "Detection blink: hold the previous command. Default 2 s.",
     ROVER),
    ("Absent · past delay",
     "on_target_lost(last_detection)",
     "Drone yaw-spins toward last side. Rover yaw-spins at search_yawspeed_slow.",
     WARN),
]
y = Inches(2.8)
for label, code, body, color in branches:
    add_box(s, Inches(0.5), y, Inches(3.8), Inches(1.05), color,
            text=label, text_size=12)
    add_text(
        s, Inches(4.5), y + Inches(0.04), Inches(8.3), Inches(0.45),
        [code], size=12, mono=True, color=INK,
    )
    add_text(
        s, Inches(4.5), y + Inches(0.5), Inches(8.3), Inches(0.55),
        [body], size=11, color=MUTED, italic=True,
    )
    y += Inches(1.25)

add_text(
    s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
    [
        "Why it matters: the legacy live_control_loop was hard-coded to MAVSDK. "
        "Extracting this loop is what let the rover adapter ride existing logic for free.",
    ],
    size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)


# Slide 6 — Drone vs Rover side by side
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Two adapters · one interface",
          "Same protocol, completely different robot-specific behaviours.")

# Drone column
add_box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.6),
        DRONE, text="MavsdkDroneAdapter", text_size=14)
drone_items = [
    ("caps", "{FORWARD, YAW, ALTITUDE}  ·  deg/s"),
    ("wire", "MAVSDK · PX4 offboard set_velocity_body"),
    ("safety", "tilt-induced retreat-from-frame-edge"),
    ("safety", "altitude clamps (min/max/down-speed)"),
    ("smoothing", "EMA + max_forward_accel slew-rate cap"),
    ("lifecycle", "arm · takeoff · land · graceful disarm"),
    ("on lost", "yaw-spin toward last-seen side · 10 °/s"),
]
y = Inches(2.05)
for label, body in drone_items:
    add_box(s, Inches(0.5), y, Inches(1.4), Inches(0.5), DRONE,
            text=label, text_size=10)
    add_text(
        s, Inches(2.0), y + Inches(0.05), Inches(4.5), Inches(0.45),
        [body], size=11, color=INK,
    )
    y += Inches(0.58)

# Rover column
add_box(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.6),
        ROVER, text="Ros2RoverAdapter", text_size=14)
rover_items = [
    ("caps", "{FORWARD, YAW}  ·  rad/s  ·  no ALTITUDE"),
    ("wire", "rclpy · /cmd_vel geometry_msgs/Twist"),
    ("safety", "bottom-edge slow-stop  (bbox_bottom ≥ 0.85)"),
    ("frame", "NED → ENU sign flip on angular.z"),
    ("smoothing", "per-axis EMA (yaw_alpha · forward_alpha)"),
    ("lifecycle", "SIGINT-safe shutdown · executor thread join"),
    ("on lost", "yaw-spin toward last-seen side · 0.3 rad/s"),
]
y = Inches(2.05)
for label, body in rover_items:
    add_box(s, Inches(6.8), y, Inches(1.4), Inches(0.5), ROVER,
            text=label, text_size=10)
    add_text(
        s, Inches(8.3), y + Inches(0.05), Inches(4.5), Inches(0.45),
        [body], size=11, color=INK,
    )
    y += Inches(0.58)

add_text(
    s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
    [
        "The controller emits the same RobotCommand into both. Units, frames, "
        "and safety overrides live entirely inside each adapter.",
    ],
    size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)


# Slide 7 — Notable fixes this branch shipped
s = prs.slides.add_slide(blank)
set_bg(s, BG_LIGHT)
add_title(s, "Notable bugs found · fixed in this branch",
          "Two that took meaningful detective work; pinned by regression tests now.")

# Bug 1 — JSON config silently shadowed
add_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.6),
        WARN, text="① JSON config silently shadowed by argparse defaults  ·  fixed in 9b2e978", text_size=14)
add_text(
    s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.7),
    [
        "add_args declared default=defaults.X (= drone dataclass default 4.0 for kp_yaw).",
        "argparse stored 4.0 on args.kp_yaw whether the user passed --yaw-gain or not.",
        "_arg returned args.kp_yaw  (non-None)  and the rover JSON's 0.05 was silently discarded.",
        "",
        "Symptom: rover at cx=0.084 emitted yaw raw = -20.95 (matches kp_yaw=4 × √27.5°) — 80× too high.",
        "Fix: argparse defaults → None.  _arg now falls through to JSON-derived defaults.",
        "Regression: test_from_args_does_not_override_json_with_cli_defaults.",
    ],
    size=12, color=INK, line_spacing=1.2,
)

# Bug 2 — on_target_lost stub
add_box(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
        WARN, text="② Rover on_target_lost was an empty-Twist stub  ·  fixed in 9fcc2af", text_size=14)
add_text(
    s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.6),
    [
        "Camera is body-fixed and forward-facing. Once the actor leaves the FoV, only physical",
        "rotation finds them again. The stub froze the rover, the 2 s grace period burned through",
        "the last commanded yaw at the DiffDrive cap (2 rad/s ≈ 230° body rotation), and the",
        "camera ended up pointing away from the actor permanently.",
        "",
        "Fix: yaw-spin at config.search_yawspeed_slow toward the side the target was last seen.",
        "Regression: three tests pin the direction logic (right, left, no-last-detection default).",
    ],
    size=12, color=INK, line_spacing=1.2,
)


# Slide 8 — Recap
s = prs.slides.add_slide(blank)
set_bg(s, BG_WHITE)
add_text(
    s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
    ["One controller. One orchestrator. Two adapters."],
    size=30, bold=True, color=INK, align=PP_ALIGN.CENTER,
)
add_text(
    s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(2.5),
    [
        "follow_api/controller.py never mentions \"drone\" or \"rover\".",
        "robot_api/orchestrator.py is the loop they both ride.",
        "Adapters carry the messy parts — units, frames, lifecycle, safety overrides.",
        "",
        "v1.1 added rover support by adding code, not by editing the controller.",
    ],
    size=17, color=MUTED, italic=True, align=PP_ALIGN.CENTER,
)
add_footer(s)


prs.save(str(OUT_PATH))
print(f"wrote {OUT_PATH}")
